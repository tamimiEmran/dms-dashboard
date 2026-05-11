"""Extract slide 3 of AI Brain Architecture_Final.pptx into reusable assets.

Outputs:
  - <out_dir>/pic_NN.<ext>  -- one file per unique embedded picture
  - <out_dir>/shapes.json   -- shape geometry + text + image references for HTML reconstruction

Coordinates are reported as fractions of the slide (so the consumer can scale to any pixel size).
"""
import hashlib
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = Path("C:/Users/Admin/Downloads/AI Brain Architecture_Final.pptx")
OUT = Path("M:/DMS/reports/dashboard/assets/leasing_taxis_arch")
OUT.mkdir(parents=True, exist_ok=True)

prs = Presentation(PPTX)
slide = prs.slides[2]  # 3rd slide (0-indexed)

slide_w = prs.slide_width
slide_h = prs.slide_height

shapes_out = []
pic_index = {}  # hash -> filename
pic_counter = 0


def shape_kind(shape):
    return {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "rect",
        MSO_SHAPE_TYPE.LINE: "line",
        MSO_SHAPE_TYPE.TEXT_BOX: "text",
        MSO_SHAPE_TYPE.PICTURE: "picture",
        MSO_SHAPE_TYPE.FREEFORM: "freeform",
        MSO_SHAPE_TYPE.GROUP: "group",
    }.get(shape.shape_type, str(shape.shape_type))


for i, shape in enumerate(slide.shapes):
    entry = {
        "i": i,
        "kind": shape_kind(shape),
        "left_frac": (shape.left or 0) / slide_w,
        "top_frac": (shape.top or 0) / slide_h,
        "w_frac": (shape.width or 0) / slide_w,
        "h_frac": (shape.height or 0) / slide_h,
        "left_emu": int(shape.left or 0),
        "top_emu": int(shape.top or 0),
        "w_emu": int(shape.width or 0),
        "h_emu": int(shape.height or 0),
    }
    if shape.has_text_frame:
        runs_per_para = []
        for para in shape.text_frame.paragraphs:
            line = para.text
            if line.strip():
                runs_per_para.append(line)
        entry["text"] = "\n".join(runs_per_para)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            blob = shape.image.blob
            ext = shape.image.ext
            h = hashlib.sha1(blob).hexdigest()[:10]
            if h not in pic_index:
                pic_counter += 1
                fname = f"pic_{pic_counter:02d}.{ext}"
                (OUT / fname).write_bytes(blob)
                pic_index[h] = fname
            entry["image"] = pic_index[h]
        except (ValueError, AttributeError) as e:
            entry["image_error"] = str(e)
    shapes_out.append(entry)

manifest = {
    "source_pptx": str(PPTX).replace("\\", "/"),
    "slide_index_zero_based": 2,
    "slide_w_emu": slide_w,
    "slide_h_emu": slide_h,
    "slide_aspect": f"{slide_w}/{slide_h}",
    "slide_size_inches": f"{slide_w/914400:.2f}x{slide_h/914400:.2f}",
    "shape_count": len(shapes_out),
    "image_count": len(pic_index),
    "shapes": shapes_out,
}

with open(OUT / "shapes.json", "w", encoding="utf-8") as fh:
    json.dump(manifest, fh, indent=2, ensure_ascii=False)

print(f"Wrote {len(shapes_out)} shapes, {len(pic_index)} unique images to {OUT}")
for h, f in pic_index.items():
    print(f"  {f}  (sha1[:10]={h})")
